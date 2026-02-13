"""
Python Tools - Python 代码执行与包管理工具

版本: v2.5.2
功能:
- 检查 Python 包是否安装
- 智能包分析：根据用户意图分析所需包
- 安装 Python 包
- 执行 Python 代码（带自动清理）
- 创建并执行临时脚本
- 区分"缺少包"和"代码错误"的情况

更新:
- 修复 RgbColor 导入错误
- 增强错误分类：区分包未安装 vs 代码导入错误
- 添加智能包检查功能
"""

import subprocess
import sys
import os
import tempfile
import importlib
import importlib.util
import re
from typing import Optional, List, Dict, Any, Tuple
from pathlib import Path

from ..schema import ToolDefinition, ToolParameter, ParameterType

import logging
logger = logging.getLogger("OpenClaw.Python")


# 项目根目录
PROJECT_ROOT = Path(__file__).parent.parent.parent


# 包名到导入名的映射
PACKAGE_TO_IMPORT = {
    "python-pptx": "pptx",
    "Pillow": "PIL",
    "opencv-python": "cv2",
    "scikit-learn": "sklearn",
    "beautifulsoup4": "bs4",
    "pyyaml": "yaml",
    "python-docx": "docx",
    "openpyxl": "openpyxl",
    "xlrd": "xlrd",
    "xlwt": "xlwt",
    "reportlab": "reportlab",
    "pypdf2": "PyPDF2",
    "pdf2image": "pdf2image",
}

# 导入名到包名的映射（反向）
IMPORT_TO_PACKAGE = {v: k for k, v in PACKAGE_TO_IMPORT.items()}
IMPORT_TO_PACKAGE.update({
    "pptx": "python-pptx",
    "PIL": "Pillow",
    "cv2": "opencv-python",
    "sklearn": "scikit-learn",
    "bs4": "beautifulsoup4",
    "yaml": "pyyaml",
    "docx": "python-docx",
})

# 用户意图到所需包的映射
TASK_TO_PACKAGES = {
    "pptx": ["python-pptx"],
    "ppt": ["python-pptx"],
    "powerpoint": ["python-pptx"],
    "word": ["python-docx"],
    "docx": ["python-docx"],
    "excel": ["openpyxl"],
    "xlsx": ["openpyxl"],
    "pdf": ["reportlab", "pypdf2"],
    "image": ["Pillow"],
    "图片": ["Pillow"],
    "图像": ["Pillow", "opencv-python"],
    "csv": [],  # 内置
    "json": [],  # 内置
    "yaml": ["pyyaml"],
    "html": ["beautifulsoup4"],
    "爬虫": ["requests", "beautifulsoup4"],
    "数据分析": ["pandas", "numpy"],
    "绘图": ["matplotlib"],
    "机器学习": ["scikit-learn", "numpy", "pandas"],
}


def check_package_installed(package_name: str) -> Dict[str, Any]:
    """
    检查 Python 包是否已安装
    
    Args:
        package_name: 包名称（如 python-pptx, pandas）
        
    Returns:
        检查结果，包含安装状态、版本等信息
    """
    logger.info(f"[Python] 检查包: {package_name}")
    
    import_name = PACKAGE_TO_IMPORT.get(package_name, package_name.replace("-", "_"))
    
    try:
        spec = importlib.util.find_spec(import_name)
        if spec is not None:
            # 尝试获取版本
            try:
                module = importlib.import_module(import_name)
                version = getattr(module, "__version__", "未知")
            except:
                version = "已安装"
            
            logger.info(f"[Python] 包 {package_name} 已安装，版本: {version}")
            return {
                "success": True,
                "installed": True,
                "package": package_name,
                "import_name": import_name,
                "version": version,
                "message": f"✅ 包 {package_name} 已安装（版本: {version}），无需重复安装。"
            }
        else:
            logger.info(f"[Python] 包 {package_name} 未安装")
            return {
                "success": True,
                "installed": False,
                "package": package_name,
                "import_name": import_name,
                "message": f"❌ 包 {package_name} 未安装。是否需要安装？请回复'是'或'安装'来确认安装。"
            }
    except Exception as e:
        logger.error(f"[Python] 检查包失败: {e}")
        return {
            "success": False,
            "error": str(e),
            "package": package_name
        }


def analyze_required_packages(task_description: str) -> Dict[str, Any]:
    """
    根据任务描述分析可能需要的 Python 包
    
    Args:
        task_description: 用户的任务描述
        
    Returns:
        分析结果，包含可能需要的包列表及其安装状态
    """
    logger.info(f"[Python] 分析任务所需包: {task_description[:50]}...")
    
    task_lower = task_description.lower()
    required_packages = []
    
    # 根据关键词匹配
    for keyword, packages in TASK_TO_PACKAGES.items():
        if keyword in task_lower:
            required_packages.extend(packages)
    
    # 去重
    required_packages = list(set(required_packages))
    
    # 检查每个包的安装状态
    package_status = []
    all_installed = True
    missing_packages = []
    
    for pkg in required_packages:
        result = check_package_installed(pkg)
        status = {
            "package": pkg,
            "installed": result.get("installed", False),
            "version": result.get("version", "")
        }
        package_status.append(status)
        
        if not status["installed"]:
            all_installed = False
            missing_packages.append(pkg)
    
    if not required_packages:
        return {
            "success": True,
            "packages_needed": [],
            "all_installed": True,
            "message": "此任务不需要额外的 Python 包。"
        }
    
    if all_installed:
        installed_list = ", ".join([f"{s['package']}({s['version']})" for s in package_status])
        return {
            "success": True,
            "packages_needed": required_packages,
            "package_status": package_status,
            "all_installed": True,
            "message": f"✅ 所需包均已安装: {installed_list}，可以直接执行任务。"
        }
    else:
        return {
            "success": True,
            "packages_needed": required_packages,
            "package_status": package_status,
            "all_installed": False,
            "missing_packages": missing_packages,
            "message": f"❌ 缺少以下包: {', '.join(missing_packages)}。是否需要安装？请回复'是'来确认。"
        }


def install_python_package(package_name: str, upgrade: bool = False) -> Dict[str, Any]:
    """
    安装 Python 包
    
    Args:
        package_name: 包名称
        upgrade: 是否升级已安装的包
        
    Returns:
        安装结果
    """
    logger.info(f"[Python] 安装包: {package_name}")
    
    try:
        # 使用 subprocess 调用 pip
        cmd = [sys.executable, "-m", "pip", "install"]
        if upgrade:
            cmd.append("--upgrade")
        cmd.append(package_name)
        
        logger.info(f"[Python] 执行命令: {' '.join(cmd)}")
        
        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=300  # 5分钟超时
        )
        
        if result.returncode == 0:
            logger.info(f"[Python] 包 {package_name} 安装成功")
            return {
                "success": True,
                "package": package_name,
                "message": f"✅ 包 {package_name} 安装成功！",
                "output": result.stdout[-500:] if len(result.stdout) > 500 else result.stdout
            }
        else:
            logger.error(f"[Python] 包安装失败: {result.stderr}")
            return {
                "success": False,
                "package": package_name,
                "error": f"安装失败: {result.stderr[-500:] if len(result.stderr) > 500 else result.stderr}"
            }
    except subprocess.TimeoutExpired:
        return {
            "success": False,
            "package": package_name,
            "error": "安装超时（超过5分钟）"
        }
    except Exception as e:
        logger.error(f"[Python] 安装异常: {e}")
        return {
            "success": False,
            "package": package_name,
            "error": str(e)
        }


def execute_python_code(
    code: str,
    description: str = "",
    working_directory: Optional[str] = None,
    timeout: int = 60
) -> Dict[str, Any]:
    """
    执行 Python 代码
    
    会创建临时脚本文件，执行后自动删除。
    
    Args:
        code: Python 代码内容
        description: 代码功能描述
        working_directory: 工作目录（相对于项目根目录）
        timeout: 执行超时时间（秒）
        
    Returns:
        执行结果
    """
    logger.info(f"[Python] 执行代码: {description or '临时脚本'}")
    
    # 确定工作目录
    if working_directory:
        work_dir = PROJECT_ROOT / working_directory
    else:
        work_dir = PROJECT_ROOT
    
    if not work_dir.exists():
        work_dir.mkdir(parents=True, exist_ok=True)
    
    # 创建临时脚本文件
    temp_script = None
    try:
        # 在工作目录创建临时文件
        fd, temp_path = tempfile.mkstemp(suffix='.py', dir=str(work_dir), prefix='_temp_script_')
        os.close(fd)
        temp_script = Path(temp_path)
        
        # 写入代码
        temp_script.write_text(code, encoding='utf-8')
        logger.info(f"[Python] 临时脚本: {temp_script}")
        
        # 执行脚本
        result = subprocess.run(
            [sys.executable, str(temp_script)],
            capture_output=True,
            text=True,
            timeout=timeout,
            cwd=str(work_dir)
        )
        
        output = result.stdout
        error = result.stderr
        
        if result.returncode == 0:
            logger.info(f"[Python] 代码执行成功")
            return {
                "success": True,
                "message": f"✅ 代码执行成功！",
                "description": description,
                "output": output[-2000:] if len(output) > 2000 else output,
                "return_code": 0
            }
        else:
            logger.error(f"[Python] 代码执行失败: {error}")
            
            # 详细分析错误类型
            error_analysis = _analyze_import_error(error)
            
            if error_analysis["error_type"] == "missing_package":
                # 真正缺少包的情况
                missing_packages = error_analysis["missing_packages"]
                return {
                    "success": False,
                    "error_type": "missing_package",
                    "error": f"执行失败，缺少以下包: {', '.join(missing_packages)}",
                    "missing_packages": missing_packages,
                    "need_install": True,
                    "message": f"❌ 执行失败，需要安装以下包：{', '.join(missing_packages)}\n是否需要安装？请回复'是'或'安装'来确认。"
                }
            elif error_analysis["error_type"] == "import_error":
                # 包已安装但导入错误（代码bug）
                import_errors = error_analysis["import_errors"]
                error_details = "\n".join([e["message"] for e in import_errors])
                return {
                    "success": False,
                    "error_type": "code_error",
                    "error": error[-1000:] if len(error) > 1000 else error,
                    "import_errors": import_errors,
                    "need_install": False,
                    "message": f"❌ 执行失败：代码中存在导入错误（非包缺失）\n\n{error_details}\n\n请修正代码中的导入语句。",
                    "suggestion": "检查导入的类名/函数名是否正确，参考包的官方文档。"
                }
            
            return {
                "success": False,
                "error_type": "runtime_error",
                "error": error[-1000:] if len(error) > 1000 else error,
                "output": output[-500:] if len(output) > 500 else output,
                "return_code": result.returncode,
                "message": f"❌ 执行失败：{error[:500]}"
            }
            
    except subprocess.TimeoutExpired:
        logger.error(f"[Python] 执行超时")
        return {
            "success": False,
            "error": f"执行超时（超过 {timeout} 秒）"
        }
    except Exception as e:
        logger.error(f"[Python] 执行异常: {e}")
        return {
            "success": False,
            "error": str(e)
        }
    finally:
        # 清理临时脚本
        if temp_script and temp_script.exists():
            try:
                temp_script.unlink()
                logger.info(f"[Python] 已清理临时脚本: {temp_script.name}")
            except Exception as e:
                logger.warning(f"[Python] 清理临时脚本失败: {e}")


def _analyze_import_error(error_message: str) -> Dict[str, Any]:
    """
    分析导入错误，区分"包未安装"和"包已安装但导入错误"
    
    Returns:
        {
            "error_type": "missing_package" | "import_error" | "other",
            "missing_packages": [...],  # 需要安装的包
            "import_errors": [...],      # 已安装但有导入问题的包
            "details": str
        }
    """
    result = {
        "error_type": "other",
        "missing_packages": [],
        "import_errors": [],
        "details": ""
    }
    
    # 1. 检测完全缺失的模块 (ModuleNotFoundError / ImportError: No module named)
    # 这种情况需要安装包
    missing_module_patterns = [
        r"ModuleNotFoundError: No module named ['\"]([^'\"]+)['\"]",
        r"ImportError: No module named ['\"]?([^\s'\"]+)['\"]?"
    ]
    
    missing_modules = []
    for pattern in missing_module_patterns:
        matches = re.findall(pattern, error_message)
        missing_modules.extend(matches)
    
    # 2. 检测"cannot import name"错误 - 这通常意味着包已安装但导入路径/名称错误
    # 这种情况不需要安装包，而是需要修复代码
    import_name_pattern = r"cannot import name ['\"]([^'\"]+)['\"] from ['\"]([^'\"]+)['\"]"
    import_name_errors = re.findall(import_name_pattern, error_message)
    
    if import_name_errors:
        # 检查这些包是否真的已安装
        for name, module in import_name_errors:
            top_module = module.split('.')[0]
            package_name = IMPORT_TO_PACKAGE.get(top_module, top_module)
            
            # 检查包是否实际安装
            check_result = check_package_installed(package_name)
            
            if check_result.get("installed"):
                # 包已安装，这是代码错误
                result["import_errors"].append({
                    "package": package_name,
                    "module": module,
                    "name": name,
                    "message": f"包 {package_name} 已安装，但 '{name}' 在 '{module}' 中不存在。这是代码错误，不是缺少包。"
                })
            else:
                # 包确实未安装
                if package_name not in result["missing_packages"]:
                    result["missing_packages"].append(package_name)
    
    # 转换缺失模块到包名
    for m in missing_modules:
        top_module = m.split('.')[0]
        package_name = IMPORT_TO_PACKAGE.get(top_module, top_module)
        if package_name not in result["missing_packages"]:
            result["missing_packages"].append(package_name)
    
    # 确定错误类型
    if result["import_errors"] and not result["missing_packages"]:
        result["error_type"] = "import_error"
        result["details"] = "包已安装，但存在导入错误（如类名或函数名不存在）。这是代码问题，需要修复代码而不是重新安装包。"
    elif result["missing_packages"]:
        result["error_type"] = "missing_package"
        result["details"] = f"缺少以下包: {', '.join(result['missing_packages'])}"
    
    return result


def _detect_missing_packages(error_message: str) -> List[str]:
    """
    从错误信息中检测缺失的包（向后兼容接口）
    只返回真正缺失的包，不包含已安装但导入错误的包
    """
    analysis = _analyze_import_error(error_message)
    return analysis.get("missing_packages", [])


def convert_file(
    input_path: str,
    output_format: str,
    output_path: Optional[str] = None,
    single_page: bool = False
) -> Dict[str, Any]:
    """
    文件格式转换（支持 txt->pptx 等）
    
    Args:
        input_path: 输入文件路径（相对于项目根目录）
        output_format: 目标格式（如 pptx, pdf, docx）
        output_path: 输出文件路径（可选，默认在同目录）
        single_page: 是否将所有内容放在一页（仅 pptx）
        
    Returns:
        转换结果
    """
    logger.info(f"[Python] 转换文件: {input_path} -> {output_format}")
    
    input_file = PROJECT_ROOT / input_path
    
    if not input_file.exists():
        return {
            "success": False,
            "error": f"输入文件不存在: {input_path}"
        }
    
    # 确定输出路径
    if output_path:
        output_file = PROJECT_ROOT / output_path
    else:
        output_file = input_file.with_suffix(f".{output_format}")
    
    # 根据转换类型确定所需包
    input_suffix = input_file.suffix.lower()
    output_format = output_format.lower()
    
    # 格式到包的映射
    FORMAT_PACKAGES = {
        "pptx": "python-pptx",
        "ppt": "python-pptx",
        "docx": "python-docx",
        "xlsx": "openpyxl",
        "pdf": "reportlab"
    }
    
    required_package = FORMAT_PACKAGES.get(output_format)
    
    if required_package:
        # 检查包是否安装
        check_result = check_package_installed(required_package)
        
        if check_result.get("installed"):
            logger.info(f"[Python] 包 {required_package} 已安装（版本: {check_result.get('version')}），直接执行转换")
        else:
            return {
                "success": False,
                "need_install": True,
                "missing_packages": [required_package],
                "package_status": "not_installed",
                "message": f"❌ 需要安装 {required_package} 包才能转换为 {output_format}。\n是否安装？请回复'是'或'安装'来确认。"
            }
    
    # 根据格式生成代码
    if output_format in ("pptx", "ppt"):
        code = _generate_txt_to_pptx_code(str(input_file), str(output_file), single_page=single_page)
    else:
        return {
            "success": False,
            "error": f"暂不支持转换为 {output_format} 格式"
        }
    
    # 执行转换
    result = execute_python_code(
        code=code,
        description=f"转换 {input_path} 为 {output_format}",
        working_directory=str(input_file.parent.relative_to(PROJECT_ROOT)) if input_file.parent != PROJECT_ROOT else None
    )
    
    if result.get("success"):
        result["output_file"] = str(output_file.relative_to(PROJECT_ROOT))
        result["message"] = f"✅ 文件转换成功！\n输出文件: {result['output_file']}"
    elif result.get("error_type") == "code_error":
        # 代码错误（如导入错误），不要误导用户认为需要安装包
        result["message"] = (
            f"❌ 转换失败：代码执行错误（非包缺失）\n\n"
            f"{result.get('error', '')[:500]}\n\n"
            f"注意：{required_package} 包已安装，这是代码逻辑错误。请检查代码或使用 execute_python 工具手动执行修复后的代码。"
        )
    
    return result


def _generate_txt_to_pptx_code(input_path: str, output_path: str, single_page: bool = False) -> str:
    """
    生成 txt 转 pptx 的 Python 代码
    
    Args:
        input_path: 输入文件路径
        output_path: 输出文件路径
        single_page: 是否将所有内容放在一页
    """
    if single_page:
        # 单页模式：所有内容放在一页
        return f'''# 自动生成的 txt -> pptx 转换脚本（单页模式）
from pptx import Presentation
from pptx.util import Inches, Pt

# 读取文本文件
input_path = r"{input_path}"
output_path = r"{output_path}"

with open(input_path, 'r', encoding='utf-8') as f:
    content = f.read()

# 创建 PPT
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 使用空白布局
blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
slide = prs.slides.add_slide(blank_layout)

# 添加文本框，覆盖大部分幻灯片区域
left = top = Inches(0.5)
width = prs.slide_width - Inches(1)
height = prs.slide_height - Inches(1)

textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
text_frame.word_wrap = True

# 设置文本
p = text_frame.paragraphs[0]
p.text = content

# 自动调整字体大小
for paragraph in text_frame.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(12)

# 保存
prs.save(output_path)
print(f"PPT 已保存到: {{output_path}}")
'''
    else:
        # 多页模式：按段落分页
        return f'''# 自动生成的 txt -> pptx 转换脚本（多页模式）
from pptx import Presentation
from pptx.util import Inches, Pt

# 读取文本文件
input_path = r"{input_path}"
output_path = r"{output_path}"

with open(input_path, 'r', encoding='utf-8') as f:
    content = f.read()

# 创建 PPT
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 按段落分割内容
paragraphs = [p.strip() for p in content.split('\\n\\n') if p.strip()]

# 如果没有明确分段，按行分割
if len(paragraphs) <= 1:
    lines = [l.strip() for l in content.split('\\n') if l.strip()]
    # 每 5 行一页
    paragraphs = []
    for i in range(0, len(lines), 5):
        paragraphs.append('\\n'.join(lines[i:i+5]))

# 创建幻灯片
for i, para in enumerate(paragraphs):
    # 添加标题幻灯片或内容幻灯片
    if i == 0:
        slide_layout = prs.slide_layouts[0]  # 标题布局
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        lines = para.split('\\n')
        title.text = lines[0] if lines else "演示文稿"
        if len(lines) > 1:
            subtitle.text = '\\n'.join(lines[1:])
    else:
        slide_layout = prs.slide_layouts[1]  # 标题和内容布局
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        lines = para.split('\\n')
        title.text = f"第 {{i}} 页"
        content_shape.text = para

# 保存
prs.save(output_path)
print(f"PPT 已保存到: {{output_path}}")
'''


# ============== 工具定义 ==============

CHECK_PACKAGE_TOOL = ToolDefinition(
    name="check_package",
    description="检查 Python 包是否已安装。在执行代码或转换文件前应先调用此工具确认环境",
    parameters=[
        ToolParameter(
            name="package_name",
            type=ParameterType.STRING,
            description="要检查的包名称（如 python-pptx, pandas, numpy）",
            required=True
        )
    ],
    handler=check_package_installed,
    category="python"
)

ANALYZE_PACKAGES_TOOL = ToolDefinition(
    name="analyze_packages",
    description="根据用户任务描述分析可能需要的 Python 包，并检查它们是否已安装。建议在执行任务前调用此工具",
    parameters=[
        ToolParameter(
            name="task_description",
            type=ParameterType.STRING,
            description="用户的任务描述（如'将txt转换成pptx'、'数据分析'）",
            required=True
        )
    ],
    handler=analyze_required_packages,
    category="python"
)

INSTALL_PACKAGE_TOOL = ToolDefinition(
    name="install_package",
    description="安装 Python 包。重要：安装前必须先用 check_package 确认包未安装。已安装的包无需重复安装",
    parameters=[
        ToolParameter(
            name="package_name",
            type=ParameterType.STRING,
            description="要安装的包名称",
            required=True
        ),
        ToolParameter(
            name="upgrade",
            type=ParameterType.BOOLEAN,
            description="是否升级已安装的包",
            required=False,
            default=False
        )
    ],
    handler=install_python_package,
    category="python"
)

EXECUTE_PYTHON_TOOL = ToolDefinition(
    name="execute_python",
    description="执行 Python 代码。会创建临时脚本，执行后自动删除",
    parameters=[
        ToolParameter(
            name="code",
            type=ParameterType.STRING,
            description="要执行的 Python 代码",
            required=True
        ),
        ToolParameter(
            name="description",
            type=ParameterType.STRING,
            description="代码功能描述",
            required=False,
            default=""
        ),
        ToolParameter(
            name="working_directory",
            type=ParameterType.STRING,
            description="工作目录（相对于项目根目录）",
            required=False,
            default=""
        ),
        ToolParameter(
            name="timeout",
            type=ParameterType.INTEGER,
            description="执行超时时间（秒），默认60秒",
            required=False,
            default=60
        )
    ],
    handler=execute_python_code,
    category="python"
)

CONVERT_FILE_TOOL = ToolDefinition(
    name="convert_file",
    description="文件格式转换，如 txt 转 pptx。会自动检查所需包是否已安装，已安装则直接执行，未安装才会提示",
    parameters=[
        ToolParameter(
            name="input_path",
            type=ParameterType.STRING,
            description="输入文件路径（相对于项目根目录）",
            required=True
        ),
        ToolParameter(
            name="output_format",
            type=ParameterType.STRING,
            description="目标格式（如 pptx, pdf, docx）",
            required=True
        ),
        ToolParameter(
            name="output_path",
            type=ParameterType.STRING,
            description="输出文件路径（可选，默认在同目录）",
            required=False
        ),
        ToolParameter(
            name="single_page",
            type=ParameterType.BOOLEAN,
            description="是否将所有内容放在一页（仅 pptx），默认 False",
            required=False,
            default=False
        )
    ],
    handler=convert_file,
    category="python"
)


# 导出所有 Python 工具
PYTHON_TOOLS = [
    CHECK_PACKAGE_TOOL,
    ANALYZE_PACKAGES_TOOL,
    INSTALL_PACKAGE_TOOL,
    EXECUTE_PYTHON_TOOL,
    CONVERT_FILE_TOOL
]


def get_python_tools() -> List[ToolDefinition]:
    """获取所有 Python 工具"""
    return PYTHON_TOOLS
